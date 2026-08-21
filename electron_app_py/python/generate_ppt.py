#!/usr/bin/env python3
"""
Electron → Python IPC 브릿지
pptx_generator5.py / verse_loader5.py 를 GUI 없이 호출한다.

stdin:  JSON { rawText, style, boldFont, outputPath, rootPath }
stdout: JSON { success, error? }
"""
import sys
import os
import json
import traceback

# ── 경로 설정 ─────────────────────────────────────────────────────────────────
HERE   = os.path.dirname(os.path.abspath(__file__))

def _resolve_root(root_from_electron):
    """참고구절_최종(3.16) 루트 경로를 결정한다."""
    # Electron이 rootPath를 전달해 주면 그것을 우선 사용
    if root_from_electron and os.path.isdir(root_from_electron):
        return root_from_electron
    # 개발 환경: .../electron_app_py/python/ → ../../ = 프로젝트 루트
    candidate = os.path.abspath(os.path.join(HERE, '..', '..'))
    return candidate

# Windows stdout/stderr도 UTF-8로 강제 설정
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')

# stdin을 바이너리로 읽어 UTF-8 디코딩 (Windows CP949 오염 방지)
raw_input = sys.stdin.buffer.read().decode('utf-8').strip()

try:
    data = json.loads(raw_input)
except json.JSONDecodeError as e:
    print(json.dumps({'success': False, 'error': f'JSON 파싱 실패: {e}'}), flush=True)
    sys.exit(1)

ROOT       = _resolve_root(data.get('rootPath', ''))
PG_DIR     = os.path.join(ROOT, 'pptx_generator')
sys.path.insert(0, ROOT)
sys.path.insert(0, PG_DIR)

# ── 모듈 임포트 ───────────────────────────────────────────────────────────────
try:
    import verse_loader5 as vl
    import pptx_generator5 as pg
    from constants import BIBLE_BOOKS, EMPHASIS_PATTERN as _EMPHASIS_PATTERN
except Exception as e:
    print(json.dumps({
        'success': False,
        'error': f'모듈 로드 실패: {e}\n경로: {PG_DIR}',
    }), flush=True)
    sys.exit(1)

# ── 66권 목록 (constants 모듈 우선) ──────────────────────────────────────────
def _get_bible_books():
    return BIBLE_BOOKS

# ── 색상 정규화: "#213337" 또는 "213337" → "213337" ──────────────────────────
def _norm_color(c):
    if isinstance(c, str) and c.startswith('#'):
        return c[1:]
    return c

def _norm_style(style):
    """style 딕셔너리 안의 color 값을 # 없는 hex 문자열로 통일한다."""
    result = {}
    for key, sub in style.items():
        result[key] = {k: (_norm_color(v) if k == 'color' else v) for k, v in sub.items()}
    return result

# ── 청킹 구절 주소 통합 헬퍼 ─────────────────────────────────────────────────
import re as _re

def _should_unify(ref_group, n):
    """단일 ref + 세미콜론 없음 + 2개 이상 슬라이드 → 청킹된 것으로 판단"""
    if n <= 1 or len(ref_group) != 1:
        return False
    # 강조 마커(굵게/밑줄) 제거 후 판단
    clean = _EMPHASIS_PATTERN.sub('', ref_group[0]).strip()
    return ';' not in clean and '\t' not in clean

def _merge_labels(first_label, last_label):
    """
    '요한계시록 21:1-3' + '요한계시록 21:4' → '요한계시록 21:1-4'
    '창세기 1:1-31' + '창세기 2:1-3' → '창세기 1:1-2:3'
    """
    fl = first_label.strip()
    ll = last_label.strip()
    if fl == ll:
        return fl

    m1 = _re.match(r'^(.*?)\s+(\d+):(\d+)(?:-(\d+)(?::(\d+))?)?$', fl)
    m2 = _re.match(r'^(.*?)\s+(\d+):(\d+)(?:-(\d+)(?::(\d+))?)?$', ll)

    if m1 and m2 and m1.group(1) == m2.group(1):
        book = m1.group(1)
        ch1  = m1.group(2)
        v1   = m1.group(3)

        if m2.group(5):
            ch2 = m2.group(4)
            v2  = m2.group(5)
        elif m2.group(4):
            ch2 = m2.group(2)
            v2  = m2.group(4)
        else:
            ch2 = m2.group(2)
            v2  = m2.group(3)

        if ch1 == ch2:
            return f"{book} {ch1}:{v1}-{v2}"
        else:
            return f"{book} {ch1}:{v1}-{ch2}:{v2}"
    return fl

def _extract_with_canonical_labels(vl, kor_data, eng_data, grouped_refs, book_abbr_map, kor_font_size=None, eng_font_size=None):
    """
    grouped_refs를 순회하며 구절을 추출한다.
    구절 수에 상관없이 각 구절별로 분리하여 추출하며,
    단일 구절의 범위가 슬라이드 용량을 초과해 청킹된 경우 해당 청크 슬라이드들의
    주소 라벨을 원래 전체 범위로 통일한다.
    """
    kor_entries, eng_entries = [], []

    for ref_group in grouped_refs:
        # ref_group 내의 세미콜론 및 다중 구절을 개별 항목으로 전개
        expanded_groups = vl._expand_ref_group(ref_group) if hasattr(vl, '_expand_ref_group') else [[r] for r in ref_group]

        for single_group in expanded_groups:
            if hasattr(vl, 'extract_passages_synchronized'):
                grp_kor, grp_eng = vl.extract_passages_synchronized(
                    kor_data, eng_data, [single_group],
                    kor_font_size=kor_font_size or 26,
                    eng_font_size=eng_font_size or 18
                )
            else:
                grp_kor = vl.extract_passages_grouped(kor_data, [single_group], font_size=kor_font_size) if kor_data is not None else []
                grp_eng = vl.extract_passages_grouped_eng(eng_data, [single_group], font_size=eng_font_size) if (eng_data is not None and hasattr(vl, 'extract_passages_grouped_eng')) else []

            n = len(grp_kor) if grp_kor else len(grp_eng)

            if _should_unify(single_group, n):
                # 청킹된 케이스: 첫·마지막 레이블로 전체 범위 계산
                if grp_kor:
                    canonical_kor = _merge_labels(grp_kor[0][0], grp_kor[-1][0])
                    grp_kor = [(canonical_kor, v, e) for _, v, e in grp_kor]
                if grp_eng:
                    canonical_eng = _merge_labels(grp_eng[0][0], grp_eng[-1][0])
                    grp_eng = [(canonical_eng, v, e) for _, v, e in grp_eng]

            kor_entries.extend(grp_kor)
            eng_entries.extend(grp_eng)

    return kor_entries, eng_entries


# ── <인용> & <교독문> 항목 분리 헬퍼 ──────────────────────────────────────────
import unicodedata as _ud

# NFC 정규화된 태그 (Windows/macOS 모두 대응)
_QUOTE_TAG          = _ud.normalize('NFC', '<인용>')
_QUOTE_PATTERN      = _re.compile(r'^<\s*(?:인용|인용구)\s*>', _re.UNICODE)
_RESPONSIVE_PATTERN = _re.compile(r'^<\s*교독문(?:\s+(.*?))?\s*>', _re.UNICODE)

_ROLE_LEADER_PAT = _re.compile(r'^(?:\[(?:인도|인도자)\]|\((?:인도|인도자)\)|<(?:인도|인도자)>|(?:인도|인도자)\s*:)\s*(.*)$', _re.UNICODE)
_ROLE_CONG_PAT   = _re.compile(r'^(?:\[(?:회중|성도|교인)\]|\((?:회중|성도|교인)\)|<(?:회중|성도|교인)>|(?:회중|성도|교인)\s*:)\s*(.*)$', _re.UNICODE)
_ROLE_ALL_PAT    = _re.compile(r'^(?:\[(?:다함께|다같이|함께)\]|\((?:다함께|다같이|함께)\)|<(?:다함께|다같이|함께)>|(?:다함께|다같이|함께)\s*:)\s*(.*)$', _re.UNICODE)

def _is_quote_body(body):
    """body가 <인용> 태그로 시작하는지 정규화 후 판단한다."""
    normalized = _ud.normalize('NFC', body)
    return bool(_QUOTE_PATTERN.match(normalized))

def _strip_quote_tag(body):
    """<인용> 태그를 제거하고 뒤 내용만 반환한다."""
    normalized = _ud.normalize('NFC', body)
    return _QUOTE_PATTERN.sub('', normalized).strip()

def _is_responsive_body(body):
    """body가 <교독문> 태그로 시작하는지 정규화 후 판단한다."""
    normalized = _ud.normalize('NFC', body.strip())
    return bool(_RESPONSIVE_PATTERN.match(normalized))

def _normalize_responsive_line(line):
    line = line.strip()
    if not line:
        return None
    m_lead = _ROLE_LEADER_PAT.match(line)
    if m_lead:
        return ('leader', f"(인도) {m_lead.group(1).strip()}")
    m_cong = _ROLE_CONG_PAT.match(line)
    if m_cong:
        return ('congregation', f"(회중) {m_cong.group(1).strip()}")
    m_all = _ROLE_ALL_PAT.match(line)
    if m_all:
        return ('all', f"(다함께) {m_all.group(1).strip()}")
    return ('plain', line)

def _parse_responsive_item(body):
    """
    <교독문> 본문을 파싱하여 (title, list_of_slide_texts) 반환.
    """
    normalized = _ud.normalize('NFC', body.strip())
    lines = normalized.splitlines()
    if not lines:
        return '교독문', []

    first_line = lines[0].strip()
    m = _RESPONSIVE_PATTERN.match(first_line)
    title_extra = ''
    content_lines = lines

    if m:
        inside = (m.group(1) or '').strip()
        outside = _RESPONSIVE_PATTERN.sub('', first_line).strip()
        title_extra = inside or outside
        content_lines = lines[1:]

    if title_extra:
        sub = _re.sub(r'^교독문\s*', '', title_extra).strip()
        title = f"교독문\n{sub}" if sub else "교독문"
    else:
        title = "교독문"

    parsed_lines = []
    for cl in content_lines:
        item = _normalize_responsive_line(cl)
        if item:
            parsed_lines.append(item)

    if not parsed_lines:
        return title, []

    slides = []
    current_group = []
    has_cong = False

    for role, text in parsed_lines:
        if role == 'all':
            if current_group:
                slides.append('\n'.join([t for _, t in current_group]))
                current_group = []
                has_cong = False
            slides.append(text)
        elif role == 'leader':
            if has_cong:
                slides.append('\n'.join([t for _, t in current_group]))
                current_group = [(role, text)]
                has_cong = False
            else:
                current_group.append((role, text))
        elif role == 'congregation':
            current_group.append((role, text))
            has_cong = True
        else:  # plain line
            if not current_group:
                current_group.append(('leader', f"(인도) {text}"))
            else:
                current_group.append((role, text))

    if current_group:
        slides.append('\n'.join([t for _, t in current_group]))

    return title, slides

def _parse_quote_content(content):
    """
    <인용> 본문에서 제목/본문 분리 및 강조 서식을 파싱한다.
    '/' 기준으로 앞은 제목(title), 뒤는 본문(body).
    본문 내 강조 표기는 인라인('단어' 굵게) 및 후미('단어' 굵게) 모두 완벽 지원.
    """
    if '/' in content:
        parts = content.split('/', 1)
        q_title = parts[0].strip()
        raw_body = parts[1].strip()
    else:
        q_title = ''
        raw_body = content.strip()

    emphases = [
        {'text': m.group(1), 'kind': 'bold' if m.group(2) == '굵게' else 'underline'}
        for m in _EMPHASIS_PATTERN.finditer(raw_body)
    ]

    if not emphases:
        return q_title, raw_body, []

    text_without_emp = _EMPHASIS_PATTERN.sub('', raw_body).strip()
    if all(emp['text'] in text_without_emp for emp in emphases):
        clean_body = text_without_emp
    else:
        clean_body = _EMPHASIS_PATTERN.sub(r'\1', raw_body).strip()

    return q_title, clean_body, emphases


def _move_slide(prs, old_index, new_index):
    """python-pptx: 슬라이드를 old_index에서 new_index로 이동."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    elem = slides[old_index]
    xml_slides.remove(elem)
    if new_index >= len(xml_slides):
        xml_slides.append(elem)
    else:
        xml_slides.insert(new_index, elem)


def _insert_black_slides(prs, group_sizes):
    """
    Presentation 객체에 group_sizes에 따라
    각 번호 항목의 마지막 슬라이드 뒤에 검은 슬라이드를 삽입한다.
    group_sizes = [n1, n2, ...] (각 번호 항목이 차지하는 슬라이드 수)
    """
    from pptx.dml.color import RGBColor as _RGB

    # 삽입 위치(0-based): 각 그룹 끝 뒤
    # 뒤에서부터 처리해 인덱스 밀림을 방지
    insert_positions = []
    cumulative = 0
    for size in group_sizes:
        cumulative += size
        insert_positions.append(cumulative)

    for pos in reversed(insert_positions):
        # 빈 레이아웃으로 슬라이드 추가
        blank_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank_layout)
        # 슬라이드 배경을 검정으로 설정
        bg = new_slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _RGB(0, 0, 0)
        # 맨 끝에 추가된 슬라이드를 원하는 위치로 이동
        last_idx = len(prs.slides) - 1
        _move_slide(prs, last_idx, pos)


def _split_items(raw_text):
    """
    번호. 로 시작하는 항목들을 순서대로 분리한다.
    반환값: [('quote', 텍스트) | ('responsive', 텍스트) | ('verse', '1. 원문내용'), ...]
    """
    items = []
    lines = raw_text.strip().splitlines()
    current_lines = []

    def _flush(buf):
        if not buf:
            return
        joined = '\n'.join(buf).strip()
        body = _re.sub(r'^\d+\.\s*', '', joined, count=1).strip()
        if _is_responsive_body(body):
            items.append(('responsive', body))
        elif _is_quote_body(body):
            items.append(('quote', _strip_quote_tag(body)))
        else:
            items.append(('verse', f'1. {body}'))

    for line in lines:
        if _re.match(r'^\d+\.', line.strip()):
            _flush(current_lines)
            current_lines = [line]
        elif not current_lines and (_is_responsive_body(line.strip()) or _is_quote_body(line.strip())):
            _flush(current_lines)
            current_lines = [line]
        else:
            current_lines.append(line)
    _flush(current_lines)
    return items


# ── 메인 로직 ─────────────────────────────────────────────────────────────────
def main():
    raw_text    = data['rawText']
    style       = _norm_style(data['style'])
    bold_font   = data.get('boldFont', '나눔스퀘어 네오 ExtraBold')
    output_path = data['outputPath']
    languages   = data.get('languages', {'kor': True, 'eng': True})
    inc_kor     = bool(languages.get('kor', True))
    inc_eng     = bool(languages.get('eng', True))

    if not inc_kor and not inc_eng:
        raise ValueError('최소 하나의 언어를 선택해야 합니다.')

    # 경로
    kor_dir       = os.path.join(ROOT, 'text_DB', '개역개정-text')
    esv_file      = os.path.join(ROOT, 'text_DB', 'ESV-text', 'ESV_cleaned.txt')
    custom_tmpl   = data.get('templatePath')
    template_path = custom_tmpl if (custom_tmpl and os.path.isfile(custom_tmpl)) else os.path.join(ROOT, 'pptx_template', 'template.pptx')

    if inc_kor and not os.path.isdir(kor_dir):
        raise FileNotFoundError(f'한글 성경 폴더 없음: {kor_dir}')
    if inc_eng and not os.path.isfile(esv_file):
        raise FileNotFoundError(f'ESV 파일 없음: {esv_file}')
    if not os.path.isfile(template_path):
        raise FileNotFoundError(f'템플릿 없음: {template_path}')

    # ── 성경 데이터 로드 ──────────────────────────────────────────────────────
    bible_books = _get_bible_books()
    kor_data    = vl.load_kor_bible(kor_dir, bible_books) if inc_kor else None
    eng_data    = vl.parse_scripture_file(esv_file) if inc_eng else None

    # ── 항목 분리: <교독문> vs <인용> vs 성경 구절 ──────────────────────────
    item_list = _split_items(raw_text)
    if not item_list:
        raise ValueError(
            '입력된 항목이 없습니다. '
            '"1. 창 1:1-3" 또는 "1. <교독문> ..." 또는 "1. <인용> 텍스트" 형식으로 입력하세요.'
        )

    book_abbr_map = getattr(vl, 'book_abbr_map', {})
    main_entries, sub_entries = [], []
    group_sizes = []  # 번호 항목별 슬라이드 수 추적

    for kind, content in item_list:
        prev_len = len(main_entries)
        if kind == 'responsive':
            # <교독문> 항목: 인도/회중 페어링 슬라이드 생성, 영문란 비움
            resp_title, resp_slides = _parse_responsive_item(content)
            for st in resp_slides:
                main_entries.append((resp_title, st, []))
                sub_entries.append(('', '', []))
        elif kind == 'quote':
            # <인용> 항목: '/'가 있으면 앞쪽은 제목란, 뒤쪽은 본문란에 삽입
            q_title, clean_body, q_emphases = _parse_quote_content(content)
            main_entries.append((q_title, clean_body, q_emphases))
            sub_entries.append(('', '', []))
        else:
            # 일반 성경 구절 처리
            grouped_refs = vl.parse_multi_refs_line(content)
            if not grouped_refs:
                continue  # 파싱 실패 항목은 건너뜀
            kor_body_size = style.get('kor_body', {}).get('size', 28)
            eng_body_size = style.get('eng_body', {}).get('size', 18)

            k, e = _extract_with_canonical_labels(
                vl, kor_data, eng_data, grouped_refs, book_abbr_map,
                kor_font_size=kor_body_size, eng_font_size=eng_body_size
            )
            # 언어 설정에 따른 엔트리 분기:
            # 1) 둘 다 선택: 메인=한글, 서브=영어
            # 2) 한국어만 선택: 메인=한글, 서브=빈값(영어 박스 비움)
            # 3) 영어만 선택: 메인=영어(원래 한글 박스에 출력), 서브=빈값(영어 박스 비움)
            if inc_kor and inc_eng:
                main_entries.extend(k)
                sub_entries.extend(e)
            elif inc_kor:
                main_entries.extend(k)
                sub_entries.extend([('', '', []) for _ in k])
            else:  # inc_eng only
                main_entries.extend(e)
                sub_entries.extend([('', '', []) for _ in e])

        added = len(main_entries) - prev_len
        if added > 0:
            group_sizes.append(added)

    if not main_entries:
        raise ValueError(
            '슬라이드를 생성할 수 없습니다. '
            '유효한 성경 구절 또는 <교독문> / <인용> 항목을 입력하세요.'
        )

    # output 폴더 생성
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # PPT 생성 (Presentation 객체를 메모리에서 받음)
    prs = pg.add_scripture_to_ppt(
        template_path,
        main_entries,
        sub_entries,
        style,
        bold_font,
        output_path,
        return_prs=True,
    )

    # 번호 항목별 검은 슬라이드 삽입 (메모리상에서 처리)
    if len(group_sizes) > 0:
        _insert_black_slides(prs, group_sizes)

    # 최종 저장 (한 번만 디스크에 씀)
    prs.save(output_path)

    print(json.dumps({'success': True}), flush=True)


if __name__ == '__main__':
    try:
        main()
    except Exception as e:
        print(json.dumps({
            'success': False,
            'error': str(e),
            'detail': traceback.format_exc(),
        }), flush=True)
        sys.exit(1)