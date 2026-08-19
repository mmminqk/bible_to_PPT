import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import copy

from pptx_generator.verse_loader2 import resource_path, absolute_path, process_text, read_files_in_directory, parse_scripture_file, extract_passages_grouped, extract_passages_grouped_eng, split_and_format_verses, parse_multi_refs_line

# 경로 설정 (상대경로 기준)
KOR_BIBLE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/text_DB/개역개정-text'
ESV_BIBLE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/text_DB/ESV-text/ESV_cleaned.txt'
TEMPLATE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/pptx_template/template.pptx'
OUTPUT_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/pptx_template/output.pptx'

SEP_NUM = 3

bible_books = [
    "창세기", "출애굽기", "레위기", "민수기", "신명기", "여호수아", "사사기", "룻기", "사무엘상", "사무엘하",
    "열왕기상", "열왕기하", "역대상", "역대하", "에스라", "느헤미야", "에스더", "욥기", "시편", "잠언",
    "전도서", "아가", "이사야", "예레미야", "예레미야애가", "에스겔", "다니엘", "호세아", "요엘", "아모스",
    "오바댜", "요나", "미가", "나훔", "하박국", "스바냐", "학개", "스가랴", "말라기", "마태복음", "마가복음",
    "누가복음", "요한복음", "사도행전", "로마서", "고린도전서", "고린도후서", "갈라디아서", "에베소서", "빌립보서", "골로새서",
    "데살로니가전서", "데살로니가후서", "디모데전서", "디모데후서", "디도서", "빌레몬서", "히브리서", "야고보서", "베드로전서", "베드로후서",
    "요한일서", "요한이서", "요한삼서", "유다서", "요한계시록"
]

DEFAULT_STYLE = {
    "kor_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 37.3, "color": "#1F3337"},
    "kor_body": {"font": "나눔스퀘어 네오 Bold", "size": 28, "color": "#1F3337"},
    "eng_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 28, "color": "#8FA79F"},
    "eng_body": {"font": "Pretendard Variable", "size": 20, "color": "#4F655E"}
}

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def duplicate_slide_with_blank_layout(prs, slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)  # 빈 레이아웃 사용
    for shape in slide.shapes:
        new_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

def add_scripture_to_ppt(template_path, verse_texts, verse_texts_eng, style, output_path="output.pptx"):
    if not os.path.exists(resource_path(template_path)):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {template_path}")
    
    prs = Presentation(template_path)
    
    # 슬라이드 수 부족하면 복제
    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])
        
    for idx, (address, verse) in enumerate(verse_texts):
        slide = prs.slides[idx]
        
        # 2번째 텍스트 상자 (인덱스 1)에 개역개정 주소 텍스트 추가
        text_shape = slide.shapes[1]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if len(address.split('\n')) > 2:
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                run = p.add_run()
                # p.text = line
                # p.font.color.rgb = hex_to_rgb(style['kor_title']['color'])
                # p.font.size = Pt(style['kor_title']['size'])
                # p.font.name = style['kor_title']['font']
                run.text = line
                run.font.color.rgb = hex_to_rgb(style['kor_title']['color'])
                run.font.size = Pt(style['kor_title']['size'])
                run.font.name = style['kor_title']['font']
            else:
                words = line.split()
                if not words:
                    continue
                for j, word in enumerate(words):
                    if i == 0 and j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = word
                    run.font.color.rgb = hex_to_rgb(style['kor_title']['color'])
                    run.font.size = Pt(style['kor_title']['size'])
                    run.font.name = style['kor_title']['font']

        # 7번째 텍스트 상자 (인덱스 6)에 개역개정 본문 텍스트 추가
        text_shape = slide.shapes[6]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            run = p.add_run()
            superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
            if len(address.split('\n')) > 2:
                sup = str(address.split('\n')[i].split(' ')[1]).translate(superscript_map)
            else:
                sup = str(line.split(' ')[0]).translate(superscript_map)
            # p.text = sup+' '+' '.join(line.split(' ')[1:])
            # p.font.color.rgb = hex_to_rgb(style['kor_body']['color'])
            # p.font.size = Pt(style['kor_body']['size'])
            # p.font.name = style['kor_body']['font']
            run.text = sup+' '+' '.join(line.split(' ')[1:])
            run.font.color.rgb = hex_to_rgb(style['kor_body']['color'])
            if len(verse) > 170:
                run.font.size = Pt(style['kor_body']['size'])*0.9
            else:
                run.font.size = Pt(style['kor_body']['size'])
            run.font.name = style['kor_body']['font']

        prs.save(output_path)

    for idx, (address, verse) in enumerate(verse_texts_eng):
        slide = prs.slides[idx]

        # 6번째 텍스트 상자 (인덱스 5)에 ESV 주소 텍스트 추가
        text_shape = slide.shapes[5]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = hex_to_rgb(style['eng_title']['color'])
            p.font.size = Pt(style['eng_title']['size'])
            p.font.name = style['eng_title']['font']


        # 8번째 텍스트 상자 (인덱스 7)에 ESV 본문 텍스트 추가
        text_shape = slide.shapes[7]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        # text_frame.text = verse
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
            sup = str(line.split(' ')[0]).translate(superscript_map)
            if len(address.split('\n')) > 2:
                sup = str(address.split('\n')[i].split(' ')[1]).translate(superscript_map)
            else:
                sup = str(line.split(' ')[0]).translate(superscript_map)
            p.text = sup+' '+' '.join(line.split(' ')[1:])
            p.font.color.rgb = hex_to_rgb(style['eng_body']['color'])
            p.font.size = Pt(style['eng_body']['size'])
            p.font.name = style['eng_body']['font']

        prs.save(output_path)

def collect_style():
    return {
        "kor_title": {
            "font": kor_title_font.get(),
            "size": float(kor_title_size.get()),
            "color": kor_title_color.get()
        },
        "kor_body": {
            "font": kor_body_font.get(),
            "size": float(kor_body_size.get()),
            "color": kor_body_color.get()
        },
        "eng_title": {
            "font": eng_title_font.get(),
            "size": float(eng_title_size.get()),
            "color": eng_title_color.get()
        },
        "eng_body": {
            "font": eng_body_font.get(),
            "size": float(eng_body_size.get()),
            "color": eng_body_color.get()
        }
    }

def on_generate_click():
    # SEP_NUM = var.get()
    raw_text = input_text.get("1.0", tk.END)
    if '–' in raw_text:
        raw_text = raw_text.replace('–', '-')
    grouped_refs = parse_multi_refs_line(raw_text)
    texts = []
    for ref in grouped_refs:
        if isinstance(ref, list):
            processed = [process_text(r) for r in ref]
        else:
            processed = process_text(ref)
        texts.append(processed)
    
    texts = read_files_in_directory(absolute_path(KOR_BIBLE_PATH))
    bible_dict = {bible_books[i]: texts[i] for i in range(len(bible_books))}
    formatted_bible = split_and_format_verses(bible_dict)
    extracted_kor = extract_passages_grouped(formatted_bible, grouped_refs)

    eng_texts = parse_scripture_file(resource_path(ESV_BIBLE_PATH))
    extracted_eng = extract_passages_grouped_eng(eng_texts, grouped_refs)
    if not extracted_kor:
        messagebox.showerror("오류", "유효한 구절을 입력하세요.")
        return
    save_path = OUTPUT_PATH
    if save_path:
        style = collect_style()
        add_scripture_to_ppt(template_path=TEMPLATE_PATH, verse_texts=extracted_kor, verse_texts_eng=extracted_eng, style=style, output_path=save_path)
        os.startfile(save_path)

# ------------------------- 실행 -------------------------

root = tk.Tk()
root.title("성경 구절 PPT 변환기")
tk.Label(root, text="구절을 입력하세요 (예: 고전 13:4-7; 시 23:1-6)").pack(pady=5)
input_text = tk.Text(root, height=15, width=60)
input_text.pack(padx=10)

# var = tk.IntVar(value=3)

# tk.Label(root, text="몇 절씩 나눌까요?").pack(pady=5)
# sep_num = tk.Spinbox(root, from_=1, to=6, increment=1, width=10, textvariable=var)
# sep_num.pack(padx=10)


style_frame = tk.LabelFrame(root, text="서식 설정 (폰트, 크기, 색상)")
style_frame.pack(padx=10, pady=10)

def add_style_row(parent, label, row, default):
    tk.Label(parent, text=label).grid(row=row, column=0, sticky="w")
    font = tk.Entry(parent, width=20)
    font.insert(0, default["font"])
    font.grid(row=row, column=1)
    size = tk.Spinbox(parent, from_=6, to=72, increment=0.5, width=5)
    size.delete(0, "end")
    size.insert(0, default["size"])
    size.grid(row=row, column=2)
    color = tk.Entry(parent, width=10)
    color.insert(0, default["color"])
    color.grid(row=row, column=3)
    return font, size, color

kor_title_font, kor_title_size, kor_title_color = add_style_row(style_frame, "한글 제목", 0, DEFAULT_STYLE["kor_title"])
kor_body_font, kor_body_size, kor_body_color = add_style_row(style_frame, "한글 본문", 1, DEFAULT_STYLE["kor_body"])
eng_title_font, eng_title_size, eng_title_color = add_style_row(style_frame, "영어 제목", 2, DEFAULT_STYLE["eng_title"])
eng_body_font, eng_body_size, eng_body_color = add_style_row(style_frame, "영어 본문", 3, DEFAULT_STYLE["eng_body"])

tk.Button(root, text="PPT로 변환", command=on_generate_click).pack(pady=10)

def main():
    global root
    root.mainloop()
