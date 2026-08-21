from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import sys
import copy
import re as _re

# 직접 실행(python pptx_generator5.py)과 패키지 실행(-m) 모두 지원
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

from constants import (
    BIBLE_BOOKS, SUPERSCRIPT_MAP,
    DEFAULT_STYLE, DEFAULT_BOLD_FONT,
    EMPHASIS_BOLD, EMPHASIS_UNDERLINE,
)
from verse_loader5 import (
    resource_path, absolute_path,
    load_kor_bible, parse_scripture_file,
    extract_passages_grouped, extract_passages_grouped_eng,
    parse_multi_refs_line,
)

# ─── 하위호환 alias ──────────────────────────────────────────────────────────
bible_books = BIBLE_BOOKS

# ─── 경로 설정 ───────────────────────────────────────────────────────────────
_BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
KOR_BIBLE_PATH = os.path.join(_BASE, 'text_DB', '개역개정-text')
ESV_BIBLE_PATH = os.path.join(_BASE, 'text_DB', 'ESV-text', 'ESV_cleaned.txt')
TEMPLATE_PATH  = os.path.join(_BASE, 'pptx_template', 'template.pptx')
OUTPUT_PATH    = os.path.join(_BASE, 'pptx_template', 'output.pptx')


# ─── 유틸리티 ────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_color):
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _apply_run_style(run, font_name, font_size, color_hex, underline=False):
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.underline = underline

def _superscript(token):
    return token.translate(SUPERSCRIPT_MAP)

def _verse_body_text(line):
    """절 번호 토큰을 제거하고 본문만 반환."""
    return ' '.join(line.split(' ')[1:])

def duplicate_slide_with_blank_layout(prs, slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    for shape in slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            copy.deepcopy(shape.element), 'p:extLst'
        )
    return new_slide


# ─── 강조 구간 분할 ──────────────────────────────────────────────────────────

def _split_by_emphases(text, emphases):
    """
    text를 emphases 목록 기준으로 분할해
    [(segment_text, kind_or_None), ...] 리스트로 반환.
    kind_or_None이 None이면 일반 텍스트, 'bold'/'underline'이면 강조 구간.

    emphases 중 text 안에 실제로 존재하는 것만 적용하며,
    앞에서부터 순서대로 처리한다.
    """
    # 전체 텍스트에서 각 강조 텍스트의 위치를 수집 후 정렬
    hits = sorted(
        ((text.find(emp['text']), emp['text'], emp['kind'])
         for emp in emphases if text.find(emp['text']) != -1),
        key=lambda x: x[0],
    )

    segments = []
    cursor = 0
    for start, emph_text, kind in hits:
        if start < cursor:
            continue  # 이미 처리된 구간이면 건너뜀
        if start > cursor:
            segments.append((text[cursor:start], None))
        segments.append((emph_text, kind))
        cursor = start + len(emph_text)

    if cursor < len(text):
        segments.append((text[cursor:], None))

    return segments if segments else [(text, None)]


# ─── run 단위 텍스트 쓰기 ───────────────────────────────────────────────────

def _write_paragraph_with_emphasis(p, text, base_font, base_size, base_color, emphases, bold_font):
    """
    단락 p에 text를 강조 구간에 따라 여러 run으로 분할해 기록.
    - 일반 구간 : base_font / base_size / base_color
    - 굵게 구간 : bold_font / base_size / base_color
    - 밑줄 구간 : base_font / base_size / base_color + underline=True
    """
    segments = _split_by_emphases(text, emphases)
    for seg_text, kind in segments:
        run = p.add_run()
        run.text = seg_text
        if kind == EMPHASIS_BOLD:
            _apply_run_style(run, bold_font, base_size, base_color)
        elif kind == EMPHASIS_UNDERLINE:
            _apply_run_style(run, base_font, base_size, base_color, underline=True)
        else:
            _apply_run_style(run, base_font, base_size, base_color)

def _write_paragraph_plain(p, text, font, size, color, use_run=True):
    """강조 없이 단락 하나를 기록 (제목 등 단순 텍스트용)."""
    if use_run:
        run = p.add_run()
        run.text = text
        _apply_run_style(run, font, size, color)
    else:
        p.text = text
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = hex_to_rgb(color)

def _set_text_lines(tf, lines, font, size, color, use_run=True):
    """tf를 초기화한 뒤 lines를 강조 없이 한 줄씩 추가."""
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _write_paragraph_plain(p, line, font, size, color, use_run=use_run)


# ─── 템플릿 도형 자동 매핑 (태그 및 도형 이름 탐색) ──────────────────────────

_TAG_PATTERNS = {
    'kor_title': [
        r'\{\{\s*(?:한글_?제목|한글_?주소|kor_?title|kor_?addr)\s*\}\}',
        r'\{\s*(?:한글_?제목|한글_?주소|kor_?title|kor_?addr)\s*\}',
        r'구절\s*주소\s*\(\s*한\s*\)',
    ],
    'kor_body': [
        r'\{\{\s*(?:한글_?본문|한글_?말씀|개역개정|kor_?body|kor_?text)\s*\}\}',
        r'\{\s*(?:한글_?본문|한글_?말씀|개역개정|kor_?body|kor_?text)\s*\}',
        r'개역개정\s*본문',
    ],
    'eng_title': [
        r'\{\{\s*(?:영어_?제목|영어_?주소|eng_?title|eng_?addr)\s*\}\}',
        r'\{\s*(?:영어_?제목|영어_?주소|eng_?title|eng_?addr)\s*\}',
        r'구절\s*주소\s*\(\s*영\s*\)',
    ],
    'eng_body': [
        r'\{\{\s*(?:영어_?본문|영어_?말씀|esv|eng_?body|eng_?text)\s*\}\}',
        r'\{\s*(?:영어_?본문|영어_?말씀|esv|eng_?body|eng_?text)\s*\}',
        r'영어\s*본문',
        r'^abc$',
    ],
}

_NAME_KEYWORDS = {
    'kor_title': ['kor_title', 'kor_addr', '한글제목', '한글주소', '한글_제목'],
    'kor_body':  ['kor_body', 'kor_text', '한글본문', '한글_본문', '개역개정'],
    'eng_title': ['eng_title', 'eng_addr', '영어제목', '영어주소', '영어_제목'],
    'eng_body':  ['eng_body', 'eng_text', '영어본문', '영어_본문', 'esv'],
}

def detect_shape_mapping(slide):
    """
    슬라이드 내의 도형들을 스캔하여
    한글 제목/본문, 영어 제목/본문 텍스트 프레임의 shape 인덱스를 자동으로 찾아 매핑한다.
    
    1단계: 텍스트 내용에서 태그 패턴({{한글제목}}, {{한글본문}} 등) 검색
    2단계: 도형 이름(Selection Pane name)에서 키워드 검색
    3단계: 미매칭 항목에 대해 레거시 템플릿(인덱스 1, 6, 5, 7) fallback
    """
    mapping = {
        'kor_title': None,
        'kor_body':  None,
        'eng_title': None,
        'eng_body':  None,
    }
    
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        raw_text = shape.text_frame.text.strip()
        name_lower = shape.name.strip().lower()
        
        # 1단계: 텍스트 태그 매칭
        for key, patterns in _TAG_PATTERNS.items():
            if mapping[key] is not None:
                continue
            for pat in patterns:
                if _re.search(pat, raw_text, _re.IGNORECASE):
                    mapping[key] = idx
                    break
        
        # 2단계: 도형 이름 매칭
        for key, keywords in _NAME_KEYWORDS.items():
            if mapping[key] is not None:
                continue
            for kw in keywords:
                if kw in name_lower:
                    mapping[key] = idx
                    break

    # 3단계: 기본 템플릿 호환용 폴백 (레거시 인덱스)
    num_shapes = len(slide.shapes)
    if mapping['kor_title'] is None and num_shapes > 1 and slide.shapes[1].has_text_frame:
        mapping['kor_title'] = 1
    if mapping['kor_body'] is None and num_shapes > 6 and slide.shapes[6].has_text_frame:
        mapping['kor_body'] = 6
    if mapping['eng_title'] is None and num_shapes > 5 and slide.shapes[5].has_text_frame:
        mapping['eng_title'] = 5
    if mapping['eng_body'] is None and num_shapes > 7 and slide.shapes[7].has_text_frame:
        mapping['eng_body'] = 7

    return mapping


# ─── 슬라이드 채우기 (공통 헬퍼) ──────────────────────────────────────────────

def _fill_slide(slide, address, verse, emphases,
                title_shape_idx, body_shape_idx,
                title_style, body_style, bold_font,
                title_use_run=True):
    """슬라이드 하나에 제목(주소)과 본문을 채운다."""
    if title_shape_idx is None and body_shape_idx is None:
        return

    addr_lines = address.split('\n') if address else []
    body_lines = verse.split('\n') if verse else []
    is_multi   = len(addr_lines) > 2

    # 제목 텍스트프레임
    if title_shape_idx is not None and title_shape_idx < len(slide.shapes):
        shape = slide.shapes[title_shape_idx]
        if shape.has_text_frame:
            tf = shape.text_frame
            if not address or not address.strip():
                tf.clear()
            elif title_use_run:
                title_words = addr_lines if is_multi else [w for ln in addr_lines for w in ln.split() if w]
                _set_text_lines(tf, title_words, use_run=True, **title_style)
            else:
                _set_text_lines(tf, addr_lines, use_run=False, **title_style)

    # 본문 텍스트프레임
    if body_shape_idx is not None and body_shape_idx < len(slide.shapes):
        shape = slide.shapes[body_shape_idx]
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            if not verse or not verse.strip():
                pass
            else:
                body_size = body_style['size']
                for i, line in enumerate(body_lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    line_stripped = line.strip()

                    # 교독문 역할 태그 확인
                    is_resp = any(line_stripped.startswith(prefix) for prefix in (
                        '(인도)', '[인도]', '<인도>',
                        '(회중)', '[회중]', '<회중>',
                        '(다함께)', '[다함께]', '<다함께>',
                        '(성도)', '[성도]', '(교인)', '[교인]',
                        '(다같이)', '[다같이]', '(함께)', '[함께]'
                    ))

                    if is_resp:
                        p.space_after = Pt(10)
                        p.line_spacing = 1.2
                    else:
                        p.space_after = Pt(3)
                        p.line_spacing = 1.15

                    line_tokens = line.split()
                    if not line_tokens:
                        full_text = ''
                    elif is_multi:
                        addr_tokens = addr_lines[i].split() if i < len(addr_lines) else []
                        sup_token = addr_tokens[1] if len(addr_tokens) > 1 else ''
                        full_text = f"{_superscript(sup_token)} {_verse_body_text(line)}" if sup_token else line
                    elif line_tokens[0].isdigit():
                        sup_token = line_tokens[0]
                        full_text = f"{_superscript(sup_token)} {_verse_body_text(line)}"
                    else:
                        full_text = line

                    # 회중 또는 다함께 단락인 경우 폰트를 bold_font로 자동 적용
                    para_font = body_style['font']
                    if any(full_text.strip().startswith(prefix) for prefix in (
                        '(회중)', '[회중]', '<회중>',
                        '(성도)', '[성도]', '(교인)', '[교인]',
                        '(다함께)', '[다함께]', '<다함께>',
                        '(다같이)', '[다같이]', '(함께)', '[함께]'
                    )):
                        para_font = bold_font

                    _write_paragraph_with_emphasis(
                        p, full_text,
                        base_font  = para_font,
                        base_size  = body_size,
                        base_color = body_style['color'],
                        emphases   = emphases,
                        bold_font  = bold_font,
                    )


# ─── PPT 생성 ────────────────────────────────────────────────────────────────

def add_scripture_to_ppt(template_path, verse_texts, verse_texts_eng, style, bold_font, output_path="output.pptx", return_prs=False):
    """
    verse_texts / verse_texts_eng : [(label, verse_text, emphases), ...]
    bold_font                     : '굵게' 서식에 사용할 폰트명
    return_prs                    : True이면 저장하지 않고 Presentation 객체 반환
    """
    if not os.path.exists(resource_path(template_path)):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {template_path}")

    prs = Presentation(template_path)
    if not prs.slides:
        raise ValueError("템플릿에 슬라이드가 없습니다.")

    # 템플릿의 첫 번째 슬라이드에서 태그 및 도형 매핑 자동 감지
    mapping = detect_shape_mapping(prs.slides[0])

    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])

    # ── 한글(개역개정 또는 메인 영역) 슬라이드 채우기 ────────────────────────
    if mapping['kor_body'] is not None or mapping['kor_title'] is not None:
        for idx, (address, verse, emphases) in enumerate(verse_texts):
            _fill_slide(
                slide           = prs.slides[idx],
                address         = address,
                verse           = verse,
                emphases        = emphases,
                title_shape_idx = mapping['kor_title'],
                body_shape_idx  = mapping['kor_body'],
                title_style     = style['kor_title'],
                body_style      = style['kor_body'],
                bold_font       = bold_font,
                title_use_run   = True,
            )

    # ── 영어(ESV) 슬라이드 채우기 (매핑되어 있는 경우) ──────────────────────
    if mapping['eng_body'] is not None or mapping['eng_title'] is not None:
        for idx in range(len(prs.slides)):
            if verse_texts_eng and idx < len(verse_texts_eng):
                address, verse, emphases = verse_texts_eng[idx]
            else:
                address, verse, emphases = '', '', []

            _fill_slide(
                slide           = prs.slides[idx],
                address         = address,
                verse           = verse,
                emphases        = emphases,
                title_shape_idx = mapping['eng_title'],
                body_shape_idx  = mapping['eng_body'],
                title_style     = style['eng_title'],
                body_style      = style['eng_body'],
                bold_font       = bold_font,
                title_use_run   = False,
            )

    if return_prs:
        return prs

    prs.save(output_path)


# ─── GUI (직접 실행 시에만 로드) ─────────────────────────────────────────────

if __name__ == '__main__':
    import tkinter as tk
    from tkinter import messagebox

    def collect_style():
        return {
            key: {"font": widgets[0].get(), "size": float(widgets[1].get()), "color": widgets[2].get()}
            for key, widgets in style_widgets.items()
        }

    def on_generate_click():
        raw_text = input_text.get("1.0", tk.END).replace('–', '-')
        inc_kor = kor_check_var.get()
        inc_eng = eng_check_var.get()

        if not inc_kor and not inc_eng:
            messagebox.showerror("오류", "최소 하나의 언어를 선택해야 합니다.")
            return

        grouped_refs = parse_multi_refs_line(raw_text)

        extracted_kor = []
        extracted_eng = []

        if inc_kor:
            formatted_bible = load_kor_bible(absolute_path(KOR_BIBLE_PATH), bible_books)
            extracted_kor = extract_passages_grouped(formatted_bible, grouped_refs)

        if inc_eng:
            eng_texts = parse_scripture_file(resource_path(ESV_BIBLE_PATH))
            extracted_eng = extract_passages_grouped_eng(eng_texts, grouped_refs)

        # 출력 구성:
        # 1) 한국어만 선택: 한글 박스에 한국어, 영어 박스 비움
        # 2) 영어만 선택: 한글 박스(메인 박스)에 영어 출력, 영어 박스 비움
        # 3) 둘 다 선택: 한글 박스에 한국어, 영어 박스에 영어 출력
        if inc_kor and inc_eng:
            main_entries = extracted_kor
            sub_entries = extracted_eng
        elif inc_kor:
            main_entries = extracted_kor
            sub_entries = [('', '', []) for _ in extracted_kor]
        else:  # inc_eng only
            main_entries = extracted_eng
            sub_entries = [('', '', []) for _ in extracted_eng]

        if not main_entries:
            messagebox.showerror("오류", "유효한 구절을 입력하세요.")
            return

        style     = collect_style()
        bold_font = bold_font_entry.get().strip() or DEFAULT_BOLD_FONT

        add_scripture_to_ppt(
            template_path    = TEMPLATE_PATH,
            verse_texts      = main_entries,
            verse_texts_eng  = sub_entries,
            style            = style,
            bold_font        = bold_font,
            output_path      = OUTPUT_PATH,
        )
        os.startfile(OUTPUT_PATH)

    # ── 윈도우 구성 ──────────────────────────────────────────────────────────
    root = tk.Tk()
    root.title("성경 구절 PPT 변환기")

    tk.Label(root, text="구절을 입력하세요 (예: 고전 13:4-7 '사랑은' 굵게)").pack(pady=5)
    input_text = tk.Text(root, height=15, width=60)
    input_text.pack(padx=10)

    # ── 언어 선택 프레임 ──────────────────────────────────────────────────────
    lang_frame = tk.LabelFrame(root, text="언어 선택")
    lang_frame.pack(padx=10, pady=(6, 0), fill='x')

    kor_check_var = tk.BooleanVar(value=True)
    eng_check_var = tk.BooleanVar(value=True)

    tk.Checkbutton(lang_frame, text="한국어 (개역개정)", variable=kor_check_var).pack(side='left', padx=10, pady=4)
    tk.Checkbutton(lang_frame, text="영어 (ESV)", variable=eng_check_var).pack(side='left', padx=10, pady=4)

    # ── 서식 설정 프레임 ──────────────────────────────────────────────────────
    style_frame = tk.LabelFrame(root, text="서식 설정 (폰트, 크기, 색상)")
    style_frame.pack(padx=10, pady=(10, 0), fill='x')

    def _add_style_row(parent, label, row, default):
        tk.Label(parent, text=label).grid(row=row, column=0, sticky='w', padx=4)
        font  = tk.Entry(parent, width=22); font.insert(0, default['font']); font.grid(row=row, column=1, padx=2)
        size  = tk.Spinbox(parent, from_=6, to=72, increment=0.5, width=5)
        size.delete(0, 'end'); size.insert(0, default['size']); size.grid(row=row, column=2, padx=2)
        color = tk.Entry(parent, width=10); color.insert(0, default['color']); color.grid(row=row, column=3, padx=2)
        return font, size, color

    style_widgets = {
        'kor_title': _add_style_row(style_frame, '한글 제목', 0, DEFAULT_STYLE['kor_title']),
        'kor_body':  _add_style_row(style_frame, '한글 본문', 1, DEFAULT_STYLE['kor_body']),
        'eng_title': _add_style_row(style_frame, '영어 제목', 2, DEFAULT_STYLE['eng_title']),
        'eng_body':  _add_style_row(style_frame, '영어 본문', 3, DEFAULT_STYLE['eng_body']),
    }

    # ── 강조 서식 설정 프레임 ─────────────────────────────────────────────────
    emph_frame = tk.LabelFrame(root, text="강조 서식 설정")
    emph_frame.pack(padx=10, pady=(6, 0), fill='x')

    tk.Label(emph_frame, text="'굵게' 폰트").grid(row=0, column=0, sticky='w', padx=4, pady=4)
    bold_font_entry = tk.Entry(emph_frame, width=30)
    bold_font_entry.insert(0, DEFAULT_BOLD_FONT)
    bold_font_entry.grid(row=0, column=1, padx=4, pady=4, sticky='w')
    tk.Label(emph_frame, text="(입력 예: 창 1:1 '태초에' 굵게  /  '빛이' 밑줄)", fg='gray').grid(
        row=1, column=0, columnspan=3, sticky='w', padx=4, pady=(0, 4)
    )

    tk.Button(root, text="PPT로 변환", command=on_generate_click).pack(pady=10)

    root.mainloop()