import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import copy

# 수정된 verse_loader2에서 필요한 함수들만 임포트
from pptx_generator.verse_loader3 import (
    resource_path, 
    absolute_path, 
    parse_multi_refs_line, 
    load_bible_data, 
    extract_passages_grouped, 
    extract_passages_grouped_eng
)

# --- 경로 설정 (PKL 파일 경로로 변경) ---
# 전처리 스크립트를 통해 생성된 pkl 파일의 경로를 지정합니다.
KOR_BIBLE_PKL = './text_DB/kor_bible.pkl'
ENG_BIBLE_PKL = './text_DB/eng_bible.pkl'
TEMPLATE_PATH = './pptx_template/template.pptx'
OUTPUT_PATH = './pptx_template/output.pptx'

DEFAULT_STYLE = {
    "kor_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 37.3, "color": "#1F3337"},
    "kor_body": {"font": "나눔스퀘어 네오 Bold", "size": 28, "color": "#1F3337"},
    "eng_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 28, "color": "#8FA79F"},
    "eng_body": {"font": "Pretendard Variable", "size": 20, "color": "#4F655E"}
}

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def duplicate_slide_with_blank_layout(prs, slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    for shape in slide.shapes:
        new_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

def add_scripture_to_ppt(template_path, verse_texts, verse_texts_eng, style, output_path="output.pptx"):
    if not os.path.exists(resource_path(template_path)):
        raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {template_path}")
    
    prs = Presentation(resource_path(template_path))
    
    # 슬라이드 수 부족하면 복제
    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])
        
    for idx, (address, verse) in enumerate(verse_texts):
        if idx >= len(prs.slides): break
        slide = prs.slides[idx]
        
        # 한글 주소 (인덱스 1)
        text_shape = slide.shapes[1]
        text_frame = text_shape.text_frame
        text_frame.clear()
        
        lines = address.split('\n')
        for i, line in enumerate(lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.color.rgb = hex_to_rgb(style['kor_title']['color'])
            run.font.size = Pt(style['kor_title']['size'])
            run.font.name = style['kor_title']['font']

        # 한글 본문 (인덱스 6)
        text_shape = slide.shapes[6]
        text_frame = text_shape.text_frame
        text_frame.clear()
        
        verse_lines = verse.split('\n')
        superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
        
        for i, line in enumerate(verse_lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run = p.add_run()
            
            # 첨자 처리 로직
            parts = line.split(' ', 1)
            sup = parts[0].translate(superscript_map)
            content = parts[1] if len(parts) > 1 else ""
            
            run.text = f"{sup} {content}"
            run.font.color.rgb = hex_to_rgb(style['kor_body']['color'])
            run.font.name = style['kor_body']['font']
            # 본문 길이에 따른 폰트 크기 조절
            run.font.size = Pt(style['kor_body']['size'] * 0.9) if len(verse) > 170 else Pt(style['kor_body']['size'])

    # 영어 본문 추가 로직
    for idx, (address, verse) in enumerate(verse_texts_eng):
        if idx >= len(prs.slides): break
        slide = prs.slides[idx]

        # 영어 주소 (인덱스 5)
        text_shape = slide.shapes[5]
        text_frame = text_shape.text_frame
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.color.rgb = hex_to_rgb(style['eng_title']['color'])
            run.font.size = Pt(style['eng_title']['size'])
            run.font.name = style['eng_title']['font']

        # 영어 본문 (인덱스 7)
        text_shape = slide.shapes[7]
        text_frame = text_shape.text_frame
        text_frame.clear()
        
        eng_lines = verse.split('\n')
        for i, line in enumerate(eng_lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            parts = line.split(' ', 1)
            sup = parts[0].translate(superscript_map)
            content = parts[1] if len(parts) > 1 else ""
            
            run = p.add_run()
            run.text = f"{sup} {content}"
            run.font.color.rgb = hex_to_rgb(style['eng_body']['color'])
            run.font.size = Pt(style['eng_body']['size'])
            run.font.name = style['eng_body']['font']

    prs.save(output_path)

def collect_style():
    return {
        "kor_title": {"font": kor_title_font.get(), "size": float(kor_title_size.get()), "color": kor_title_color.get()},
        "kor_body": {"font": kor_body_font.get(), "size": float(kor_body_size.get()), "color": kor_body_color.get()},
        "eng_title": {"font": eng_title_font.get(), "size": float(eng_title_size.get()), "color": eng_title_color.get()},
        "eng_body": {"font": eng_body_font.get(), "size": float(eng_body_size.get()), "color": eng_body_color.get()}
    }

def on_generate_click():
    raw_text = input_text.get("1.0", tk.END).strip()
    if not raw_text:
        messagebox.showwarning("경고", "구절을 입력하세요.")
        return
        
    if '–' in raw_text:
        raw_text = raw_text.replace('–', '-')
    
    try:
        # 1. 입력 텍스트 파싱
        grouped_refs = parse_multi_refs_line(raw_text)
        
        # 2. PKL 데이터 로드 (verse_loader2의 함수 사용)
        kor_bible_data = load_bible_data(KOR_BIBLE_PKL)
        eng_bible_data = load_bible_data(ENG_BIBLE_PKL)

        # 3. 구절 추출
        extracted_kor = extract_passages_grouped(kor_bible_data, grouped_refs)
        extracted_eng = extract_passages_grouped_eng(eng_bible_data, grouped_refs)

        if not extracted_kor:
            messagebox.showerror("오류", "유효한 구절을 찾을 수 없습니다.")
            return

        # 4. PPT 생성
        style = collect_style()
        save_path = absolute_path(OUTPUT_PATH)
        add_scripture_to_ppt(TEMPLATE_PATH, extracted_kor, extracted_eng, style, save_path)
        
        # 5. 완료 후 열기
        os.startfile(save_path)
        # messagebox.showinfo("완료", "PPT 생성이 완료되었습니다.")
        
    except Exception as e:
        messagebox.showerror("오류 발생", str(e))

# ------------------------- GUI 설정 -------------------------

root = tk.Tk()
root.title("성경 구절 PPT 변환기 (PKL 버전)")

tk.Label(root, text="구절을 입력하세요 (예: 1. 창 1:1-3    출 3:1-2)").pack(pady=5)
input_text = tk.Text(root, height=12, width=60)
input_text.pack(padx=10)

style_frame = tk.LabelFrame(root, text="서식 설정")
style_frame.pack(padx=10, pady=10, fill="x")

def add_style_row(parent, label, row, default):
    tk.Label(parent, text=label).grid(row=row, column=0, padx=5, sticky="w")
    font = tk.Entry(parent, width=15)
    font.insert(0, default["font"])
    font.grid(row=row, column=1, padx=5)
    size = tk.Spinbox(parent, from_=10, to=100, increment=0.5, width=5)
    size.delete(0, "end")
    size.insert(0, default["size"])
    size.grid(row=row, column=2, padx=5)
    color = tk.Entry(parent, width=10)
    color.insert(0, default["color"])
    color.grid(row=row, column=3, padx=5)
    return font, size, color

kor_title_font, kor_title_size, kor_title_color = add_style_row(style_frame, "한글 제목", 0, DEFAULT_STYLE["kor_title"])
kor_body_font, kor_body_size, kor_body_color = add_style_row(style_frame, "한글 본문", 1, DEFAULT_STYLE["kor_body"])
eng_title_font, eng_title_size, eng_title_color = add_style_row(style_frame, "영어 제목", 2, DEFAULT_STYLE["eng_title"])
eng_body_font, eng_body_size, eng_body_color = add_style_row(style_frame, "영어 본문", 3, DEFAULT_STYLE["eng_body"])

tk.Button(root, text="PPT 생성 및 열기", command=on_generate_click, height=2, bg="#e1e1e1").pack(pady=10)

def main():
    """외부 파일에서 호출할 수 있도록 main 함수를 정의합니다."""
    global root
    root.mainloop()

if __name__ == "__main__":
    # 직접 이 파일을 실행할 때도 main()이 호출되도록 합니다.
    main()