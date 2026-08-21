"""
예배 통합 슬라이드 생성 및 슬롯 치환 모듈 (pptx_merger.py)

주일 예배(sunday_template.pptx) 및 수요 예배(wednesday_template.pptx) 등의 템플릿에서
{{태그명}}이 포함된 슬라이드를 감지하고, 해당 위치에 찬양 PPT 슬라이드 및 성경 구절 슬라이드를
치환/삽입하여 단 하나의 통합 예배 PPT를 생성한다.

역순 인플레이스 치환(Reverse In-Place Replacement) 방식을 사용하여
태그 이외의 모든 슬라이드(배경 이미지, 안내 문구, 날짜, 검은 화면 등)는 100% 원본 그대로 보존된다.
"""

import os
import io
import re
import copy
from pptx import Presentation

# 태그 매칭 정규식: {{태그명}}
TAG_PATTERN = re.compile(r'\{\{\s*([^{}]+?)\s*\}\}')
EMBED_ATTR = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'


def get_slide_tags(slide):
    """
    슬라이드 내의 모든 텍스트프레임에서 {{태그명}}을 추출하여 리스트로 반환.
    태그가 없으면 빈 리스트를 반환한다.
    """
    tags = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            found = TAG_PATTERN.findall(text)
            for t in found:
                tag_name = t.strip()
                if tag_name not in tags:
                    tags.append(tag_name)
    return tags


def scan_template_slots(template_path):
    """
    템플릿 PPTX의 모든 슬라이드를 스캔하여 태그 슬롯 정보를 반환한다.
    반환값: list of dict [
        {'index': 0, 'tags': ['시작찬양'], 'is_slot': True},
        {'index': 1, 'tags': [], 'is_slot': False},
        ...
    ]
    """
    if not os.path.isfile(template_path):
        raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {template_path}")

    prs = Presentation(template_path)
    slots = []
    for idx, slide in enumerate(prs.slides):
        tags = get_slide_tags(slide)
        slots.append({
            'index': idx,
            'tags': tags,
            'is_slot': len(tags) > 0,
        })
    return slots


def is_scripture_tag(tag_name):
    """태그가 성경 구절 / 말씀 슬롯인지 확인"""
    normalized = tag_name.replace(' ', '')
    return any(k in normalized for k in ['말씀참고구절', '말씀', '성경', '성경구절', '본문', '설교'])


def _insert_slide_into_prs(target_prs, src_slide, insert_index):
    """
    src_slide의 도형, 텍스트, 이미지, 배경을 target_prs에 완벽히 복제하여 삽입한다.
    외부 Presentation 간 파트 충돌 없이 깨끗하게 복사한다.
    """
    blank_layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(blank_layout)

    # 1. 이미지 관계(Relationship) 복제 및 rId 매핑
    rid_map = {}
    for rel_id, rel in src_slide.part.rels.items():
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            rid_map[rel_id] = new_rid
        elif "image" in rel.reltype.lower():
            try:
                img_blob = rel.target_part.blob
                img_part = target_prs.part.get_or_add_image_part(io.BytesIO(img_blob))
                new_rid = new_slide.part.relate_to(img_part, rel.reltype)
                rid_map[rel_id] = new_rid
            except Exception:
                pass

    # 2. Shape 트리 복제 (필요시 r:embed ID 치환)
    for shape in src_slide.shapes:
        shape_elem = copy.deepcopy(shape.element)
        if rid_map:
            for blip in shape_elem.xpath('.//a:blip'):
                embed_rid = blip.get(EMBED_ATTR)
                if embed_rid in rid_map:
                    blip.set(EMBED_ATTR, rid_map[embed_rid])
        new_slide.shapes._spTree.insert_element_before(shape_elem, 'p:extLst')

    # 3. 배경 복제
    src_bg = src_slide._element.xpath('./p:cSld/p:bg')
    if src_bg:
        bg_elem = copy.deepcopy(src_bg[0])
        if rid_map:
            for blip in bg_elem.xpath('.//a:blip'):
                embed_rid = blip.get(EMBED_ATTR)
                if embed_rid in rid_map:
                    blip.set(EMBED_ATTR, rid_map[embed_rid])
        target_bg = new_slide._element.xpath('./p:cSld/p:bg')
        for tb in target_bg:
            tb.getparent().remove(tb)
        new_slide._element.cSld.insert(0, bg_elem)

    # 4. 맨 뒤에 추가된 슬라이드를 원하는 insert_index 위치로 이동
    xml_slides = target_prs.slides._sldIdLst
    slides = list(xml_slides)
    elem = slides[-1]
    xml_slides.remove(elem)
    if insert_index >= len(xml_slides):
        xml_slides.append(elem)
    else:
        xml_slides.insert(insert_index, elem)

    return new_slide


def _delete_slide(prs, slide_index):
    """프레젠테이션에서 특정 인덱스의 슬라이드를 삭제한다."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if 0 <= slide_index < len(slides):
        xml_slides.remove(slides[slide_index])


def merge_worship_ppt(
    template_path,
    scripture_prs=None,
    external_slots=None,
    output_path="output.pptx"
):
    """
    템플릿 PPT를 직접 로드하여 역순 인플레이스 치환으로 예배 통합 PPT를 완성한다.

    - template_path: 'pptx_template/sunday_template.pptx' 등
    - scripture_prs: 생성된 성경 구절 Presentation 객체 (또는 None)
    - external_slots: dict, 예: {'시작찬양': 'path/to/song1.pptx', '예배찬양': 'path/to/song2.pptx', ...}
    - output_path: 최종 저장될 .pptx 경로
    """
    if not os.path.isfile(template_path):
        raise FileNotFoundError(f"템플릿 파일이 존재하지 않습니다: {template_path}")

    prs = Presentation(template_path)
    if not prs.slides:
        raise ValueError("템플릿에 슬라이드가 없습니다.")

    external_slots = external_slots or {}

    # 뒤에서부터 순회하여 인덱스 변동의 영향을 받지 않도록 처리
    num_slides = len(prs.slides)
    for slide_idx in reversed(range(num_slides)):
        tmpl_slide = prs.slides[slide_idx]
        tags = get_slide_tags(tmpl_slide)

        if not tags:
            # 일반 고정 슬라이드는 100% 그대로 유지
            continue

        # 치환할 슬라이드 목록 준비
        replacement_slides = []

        for tag in tags:
            tag_clean = tag.strip()

            if is_scripture_tag(tag_clean):
                if scripture_prs and scripture_prs.slides:
                    replacement_slides.extend(list(scripture_prs.slides))
            else:
                matching_file = None
                for k, v in external_slots.items():
                    k_clean = k.replace('{', '').replace('}', '').strip()
                    if k_clean == tag_clean or k_clean.replace(' ', '') == tag_clean.replace(' ', ''):
                        matching_file = v
                        break

                if matching_file and os.path.isfile(matching_file):
                    try:
                        ext_prs = Presentation(matching_file)
                        replacement_slides.extend(list(ext_prs.slides))
                    except Exception as e:
                        print(f"Warning: 외부 파일 로드 실패 ({matching_file}): {e}")

        # replacement_slides를 slide_idx 위치에 순서대로 삽입
        for offset, src_s in enumerate(replacement_slides):
            _insert_slide_into_prs(prs, src_s, slide_idx + offset)

        # 원래의 {{태그}} 플레이스홀더 슬라이드 삭제
        original_slide_index = slide_idx + len(replacement_slides)
        _delete_slide(prs, original_slide_index)

    # 최종 디렉터리 생성 및 저장
    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    prs.save(output_path)
    return output_path
