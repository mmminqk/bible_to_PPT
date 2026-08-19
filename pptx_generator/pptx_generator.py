import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import copy

from pptx_generator.verse_loader import resource_path, absolute_path, read_files_in_directory, parse_scripture_file, extract_passages_grouped, extract_passages_grouped_eng, split_and_format_verses, parse_multi_refs_line

# ------------------------- 경로 설정 -------------------------
# 개역개정 절대경로
KOR_BIBLE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/text_DB/개역개정-text'
# ESV 절대경로
ESV_BIBLE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/text_DB/ESV-text/ESV_cleaned.txt'
# PPTX 템플릿 절대경로
TEMPLATE_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/pptx_template/template.pptx'
# PPTX 출력 절대경로
OUTPUT_PATH = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) + '/pptx_template/output.pptx'

bible_books = [
    "창세기", "출애굽기", "레위기", "민수기", "신명기", "여호수아", "사사기", "룻기", "사무엘상", "사무엘하",
    "열왕기상", "열왕기하", "역대상", "역대하", "에스라", "느헤미야", "에스더", "욥기", "시편", "잠언",
    "전도서", "아가", "이사야", "예레미야", "예레미야애가", "에스겔", "다니엘", "호세아", "요엘", "아모스",
    "오바댜", "요나", "미가", "나훔", "하박국", "스바냐", "학개", "스가랴", "말라기", "마태복음", "마가복음",
    "누가복음", "요한복음", "사도행전", "로마서", "고린도전서", "고린도후서", "갈라디아서", "에베소서", "빌립보서", "골로새서",
    "데살로니가전서", "데살로니가후서", "디모데전서", "디모데후서", "디도서", "빌레몬서", "히브리서", "야고보서", "베드로전서", "베드로후서",
    "요한일서", "요한이서", "요한삼서", "유다서", "요한계시록"
]



def duplicate_slide_with_blank_layout(prs, slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)  # 빈 레이아웃 사용
    for shape in slide.shapes:
        new_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

# PPTX 파일에 성경 구절 입력
def add_scripture_to_ppt(template_path, verse_texts, verse_texts_eng, output_path="output.pptx"):
    if not os.path.exists(resource_path(template_path)):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {template_path}")
    
    prs = Presentation(template_path)

    # 슬라이드 수 부족하면 복제
    while len(prs.slides) < len(verse_texts):
        duplicate_slide_with_blank_layout(prs, prs.slides[-1])

    for idx, (address, verse) in enumerate(verse_texts):
        slide = prs.slides[idx]
        
        # 2번째 텍스트 상자 (인덱스 1)에 개역개정 주소 텍스트 추가
        text_shape = slide.shapes[1]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if len(address.split('\n')) > 2:
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line
                p.font.color.rgb = RGBColor(31, 51, 55)
                p.font.size = Pt(37.3)
                p.font.name = '나눔스퀘어 네오 ExtraBold'
            else:
                words = line.split()
                if not words:
                    continue
                for j, word in enumerate(words):
                    if i == 0 and j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = word
                    p.font.color.rgb = RGBColor(31, 51, 55)
                    p.font.size = Pt(37.3)
                    p.font.name = '나눔스퀘어 네오 ExtraBold'

        # 7번째 텍스트 상자 (인덱스 6)에 개역개정 본문 텍스트 추가
        text_shape = slide.shapes[6]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
            if len(address.split('\n')) > 2:
                sup = str(address.split('\n')[i].split(' ')[1]).translate(superscript_map)
            else:
                sup = str(line.split(' ')[0]).translate(superscript_map)
            p.text = sup+' '+' '.join(line.split(' ')[1:])
            p.font.color.rgb = RGBColor(31, 51, 55)
            p.font.size = Pt(28)
            p.font.name = '나눔스퀘어 네오 Bold'

        prs.save(output_path)

    for idx, (address, verse) in enumerate(verse_texts_eng):
        slide = prs.slides[idx]

        # 6번째 텍스트 상자 (인덱스 5)에 ESV 주소 텍스트 추가
        text_shape = slide.shapes[5]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        for i, line in enumerate(address.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.color.rgb = RGBColor(143, 167, 159)
            p.font.size = Pt(28)
            p.font.name = '나눔스퀘어 네오 ExtraBold'


        # 8번째 텍스트 상자 (인덱스 7)에 ESV 본문 텍스트 추가
        text_shape = slide.shapes[7]
        text_frame = text_shape.text_frame # type: ignore
        text_frame.clear()
        # text_frame.text = verse
        for i, line in enumerate(verse.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
            sup = str(line.split(' ')[0]).translate(superscript_map)
            if len(address.split('\n')) > 2:
                sup = str(address.split('\n')[i].split(' ')[1]).translate(superscript_map)
            else:
                sup = str(line.split(' ')[0]).translate(superscript_map)
            p.text = sup+' '+' '.join(line.split(' ')[1:])
            p.font.color.rgb = RGBColor(79, 101, 94)
            p.font.size = Pt(20)
            p.font.name = 'Pretendard Variable'
        

        prs.save(output_path)

# ------------------------- GUI 설정 -------------------------
def on_generate_click():
    raw_text = input_text.get("1.0", tk.END)
    if '–' in raw_text:
        raw_text = raw_text.replace('–', '-')
    grouped_refs = parse_multi_refs_line(raw_text) # 구절 목록을 리스트화
    
    texts = read_files_in_directory(absolute_path(KOR_BIBLE_PATH)) # 개역개정 성경 파일 읽기
    bible_dict = {bible_books[i]: texts[i] for i in range(len(bible_books))} # 성경 책 이름과 내용을 딕셔너리로 매핑
    formatted_bible = split_and_format_verses(bible_dict) # 책 별로 구분된 딕셔너리로 변환
    extracted_kor = extract_passages_grouped(formatted_bible, grouped_refs) # 개역개정 구절 추출
    
    parsed = parse_scripture_file(resource_path(ESV_BIBLE_PATH)) # ESV 성경 파일 읽기
    extracted_eng = extract_passages_grouped_eng(parsed, grouped_refs) # ESV 구절 추출

    if not extracted_kor:
        messagebox.showerror("오류", "유효한 구절을 입력하세요.")
        return
    # save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")]) # 저장 경로 선택 가능
    save_path = OUTPUT_PATH
    # PPTX 파일 저장 후 자동 실행
    if save_path:
        add_scripture_to_ppt(template_path=TEMPLATE_PATH, verse_texts=extracted_kor, verse_texts_eng=extracted_eng, output_path=save_path)
        os.startfile(save_path)
        # messagebox.showinfo("완료", f"PPTX 파일이 저장되었습니다: {save_path}")

# ------------------------- 실행 -------------------------

root = tk.Tk()
root.title("성경 구절 PPTX 변환기")
tk.Label(root, text=(
    "구절 입력 양식\n"
    "1) 각 구절은 '숫자. 책 장:절' 형식으로 입력합니다.\n"
    "2) 한 슬라이드에 여러 구절을 넣고 싶다면, 각 구절을 세미콜론(;)으로 구분합니다.\n"
    "3) 이어진 구절로 이루어져 있다면 '책 장:절-절' 형식으로 입력합니다.\n"
    "4) 책 제목은 약어로 입력합니다.\n\n"
    "5) 인용구는 '<인용구> 내용' 형식으로 입력합니다.\n\n"
    "예)\n"
    "1. 잠 1:8; 잠 31:2\n"
    "2. 마 5:7-9\n"
    "3. <인용구> \"동해물과 백두산이 마르고 닳도록\"\n"
)).pack(pady=5)
input_text = tk.Text(root, height=15, width=60)
input_text.pack(padx=10)
tk.Button(root, text="PPTX로 변환", command=on_generate_click).pack(pady=10)
# root.mainloop()

def main():
    global root
    root.mainloop()