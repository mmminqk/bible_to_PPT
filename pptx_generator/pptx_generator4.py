import tkinter as tk
from tkinter import messagebox
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import copy

from sympy import root

from pptx_generator.verse_loader4 import (
    resource_path, load_bible_data, parse_multi_refs_line,
    extract_passages_grouped, extract_passages_grouped_eng
)

KOR_BIBLE_PKL = './text_DB/kor_bible.pkl'
ENG_BIBLE_PKL = './text_DB/eng_bible.pkl'
TEMPLATE_PATH = './pptx_template/template.pptx'
OUTPUT_PATH = './pptx_template/output.pptx'

# 기본값에 'bold_font' 항목 추가
DEFAULT_STYLE = {
    "kor_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 37.3, "color": "#1F3337"},
    "kor_body": {"font": "나눔스퀘어 네오 Bold", "size": 28, "color": "#1F3337"},
    "eng_title": {"font": "나눔스퀘어 네오 ExtraBold", "size": 28, "color": "#8FA79F"},
    "eng_body": {"font": "Pretendard Variable", "size": 20, "color": "#4F655E"},
    "special": {"bold_font": "나눔스퀘어 네오 ExtraBold"} # 굵게 입력 시 적용할 폰트
}

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def apply_formatted_text(paragraph, text, fmt_info, style_key, style):
    superscript_map = str.maketrans("0123456789:", "⁰¹²³⁴⁵⁶⁷⁸⁹˸")
    parts = text.split(' ', 1)
    sup = parts[0].translate(superscript_map)
    content = parts[1] if len(parts) > 1 else ""

    keyword = fmt_info["keyword"]
    action = fmt_info["action"]
    font_override = fmt_info["font_override"]

    # 1. 절 번호 추가
    run_sup = paragraph.add_run()
    run_sup.text = f"{sup} "
    
    # 2. 본문 처리
    if keyword and keyword in content:
        before, after = content.split(keyword, 1)
        
        # 키워드 앞
        run_pre = paragraph.add_run()
        run_pre.text = before
        
        # 키워드 (핵심 수정 부분)
        run_key = paragraph.add_run()
        run_key.text = keyword
        
        # 서식 우선순위 적용
        if action == "밑줄":
            run_key.font.underline = True
        elif action == "굵게":
            # 폰트가 굵기를 지원하지 않으므로 설정된 '굵은 폰트'로 교체
            run_key.font.name = style['special']['bold_font']
        
        if font_override: # 특정 폰트명 직접 입력 시
            run_key.font.name = font_override
            
        # 키워드 뒤
        run_post = paragraph.add_run()
        run_post.text = after
    else:
        run_main = paragraph.add_run()
        run_main.text = content

    # 3. 공통 스타일 (폰트명 미지정된 run들만 기본 폰트 적용)
    for run in paragraph.runs:
        run.font.color.rgb = hex_to_rgb(style[style_key]['color'])
        run.font.size = Pt(style[style_key]['size'])
        if not run.font.name:
            run.font.name = style[style_key]['font']

def add_scripture_to_ppt(template_path, kor_results, eng_results, style, output_path):
    prs = Presentation(resource_path(template_path))
    
    while len(prs.slides) < len(kor_results):
        blank_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank_layout)
        for shape in prs.slides[0].shapes:
            new_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), 'p:extLst')

    # --- 한글 슬라이드 적용 ---
    for idx, (addr, verse, fmt) in enumerate(kor_results):
        slide = prs.slides[idx]
        
        # 한글 주소 (인덱스 1)
        slide.shapes[1].text_frame.clear()
        p_addr = slide.shapes[1].text_frame.paragraphs[0]
        run_addr = p_addr.add_run()
        run_addr.text = addr.strip()
        
        # --- 이 부분 수정: 한글 제목 폰트 및 사이즈 적용 ---
        run_addr.font.name = style['kor_title']['font']
        run_addr.font.size = Pt(style['kor_title']['size'])
        run_addr.font.color.rgb = hex_to_rgb(style['kor_title']['color'])
        
        # 한글 본문 (인덱스 6)
        slide.shapes[6].text_frame.clear()
        for line in verse.split('\n'):
            p = slide.shapes[6].text_frame.add_paragraph()
            apply_formatted_text(p, line, fmt, 'kor_body', style)

    # --- 영어 슬라이드 적용 ---
    for idx, (addr, verse, fmt) in enumerate(eng_results):
        if idx >= len(prs.slides): break
        slide = prs.slides[idx]
        
        # 영어 주소 (인덱스 5)
        slide.shapes[5].text_frame.clear()
        p_eng_addr = slide.shapes[5].text_frame.paragraphs[0]
        run_e_addr = p_eng_addr.add_run()
        run_e_addr.text = addr.strip()
        
        # --- 이 부분 수정: 영어 제목 폰트 및 사이즈 적용 ---
        run_e_addr.font.name = style['eng_title']['font']
        run_e_addr.font.size = Pt(style['eng_title']['size'])
        run_e_addr.font.color.rgb = hex_to_rgb(style['eng_title']['color'])
        
        # 영어 본문 (인덱스 7)
        slide.shapes[7].text_frame.clear()
        for line in verse.split('\n'):
            p = slide.shapes[7].text_frame.add_paragraph()
            apply_formatted_text(p, line, fmt, 'eng_body', style)

    prs.save(output_path)

# --- GUI ---

def on_generate_click():
    raw_text = input_text.get("1.0", tk.END).strip()
    if not raw_text: return
    
    try:
        grouped_refs = parse_multi_refs_line(raw_text.replace('–', '-'))
        kor_data = load_bible_data(KOR_BIBLE_PKL)
        eng_data = load_bible_data(ENG_BIBLE_PKL)
        
        extracted_kor = extract_passages_grouped(kor_data, grouped_refs)
        extracted_eng = extract_passages_grouped_eng(eng_data, grouped_refs)
        
        style = {
            "kor_title": {"font": kor_title_f.get(), "size": float(kor_title_s.get()), "color": kor_title_c.get()},
            "kor_body": {"font": kor_body_f.get(), "size": float(kor_body_s.get()), "color": kor_body_c.get()},
            "eng_title": {"font": eng_title_f.get(), "size": float(eng_title_s.get()), "color": eng_title_c.get()},
            "eng_body": {"font": eng_body_f.get(), "size": float(eng_body_s.get()), "color": eng_body_c.get()},
            "special": {"bold_font": bold_font_entry.get()} # GUI에서 가져옴
        }
        
        add_scripture_to_ppt(TEMPLATE_PATH, extracted_kor, extracted_eng, style, OUTPUT_PATH)
        os.startfile(os.path.abspath(OUTPUT_PATH))
    except Exception as e:
        messagebox.showerror("오류", str(e))

def main():
    global input_text, kor_title_f, kor_title_s, kor_title_c, kor_body_f, kor_body_s, kor_body_c, eng_title_f, eng_title_s, eng_title_c, eng_body_f, eng_body_s, eng_body_c, bold_font_entry
    
    root = tk.Tk()
    root.title("성경 구절 PPT 변환기")
    tk.Label(root, text="구절을 입력하세요 (예: 고전 13:4-7; 시 23:1-6)").pack(pady=5)
    input_text = tk.Text(root, height=15, width=60)
    input_text.pack(padx=10)

    # 스타일 설정 영역
    f_style = tk.LabelFrame(root, text="서식 설정"); f_style.pack(padx=10, pady=5, fill="x")
    
    def add_row(parent, label, default, row):
        tk.Label(parent, text=label).grid(row=row, column=0, padx=5)
        f = tk.Entry(parent, width=15); f.insert(0, default["font"]); f.grid(row=row, column=1)
        s = tk.Entry(parent, width=5); s.insert(0, str(default["size"])); s.grid(row=row, column=2)
        c = tk.Entry(parent, width=8); c.insert(0, default["color"]); c.grid(row=row, column=3)
        return f, s, c

    kor_title_f, kor_title_s, kor_title_c = add_row(f_style, "한글 제목", DEFAULT_STYLE["kor_title"], 0)
    kor_body_f, kor_body_s, kor_body_c = add_row(f_style, "한글 본문", DEFAULT_STYLE["kor_body"], 1)
    eng_title_f, eng_title_s, eng_title_c = add_row(f_style, "영어 제목", DEFAULT_STYLE["eng_title"], 2)
    eng_body_f, eng_body_s, eng_body_c = add_row(f_style, "영어 본문", DEFAULT_STYLE["eng_body"], 3)

    # 대체 굵은 폰트 설정 행 추가
    tk.Label(f_style, text="대체 굵은 폰트").grid(row=4, column=0, padx=5, pady=5)
    bold_font_entry = tk.Entry(f_style, width=15)
    bold_font_entry.insert(0, DEFAULT_STYLE["special"]["bold_font"])
    bold_font_entry.grid(row=4, column=1, columnspan=3, sticky="w", padx=5)

    tk.Button(root, text="PPT 생성", command=on_generate_click, height=2, width=20, bg="#eee").pack(pady=10)
    root.mainloop()

if __name__ == "__main__":
    main()